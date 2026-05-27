# Motion Storyboard Generator
# Run locally with: python -m streamlit run app.py

import tempfile
from difflib import SequenceMatcher
from pathlib import Path
from typing import Dict, List, Tuple

import cv2
import numpy as np
import pytesseract
import streamlit as st
from pptx import Presentation
from pptx.util import Inches


FrameResult = Tuple[str, float, str]


OCR_LANGUAGES = {
    "English": "eng",
    "Korean": "kor",
    "Japanese": "jpn",
    "Spanish": "spa",
    "Portuguese": "por",
    "French": "fra",
    "German": "deu",
    "Italian": "ita",
    "Dutch": "nld",
    "Chinese Simplified": "chi_sim",
    "Chinese Traditional": "chi_tra",
    "Arabic": "ara",
    "Hindi": "hin",
    "English + Korean": "eng+kor",
    "English + Japanese": "eng+jpn",
    "English + Spanish": "eng+spa",
    "English + Portuguese": "eng+por",
    "English + French": "eng+fra",
    "English + German": "eng+deu",
    "English + Chinese Simplified": "eng+chi_sim",
}


def clean_text(text: str) -> str:
    return " ".join(text.replace("\n", " ").split()).strip()


def get_video_duration(video_path: str) -> float:
    capture = cv2.VideoCapture(video_path)
    fps = capture.get(cv2.CAP_PROP_FPS)
    total_frames = capture.get(cv2.CAP_PROP_FRAME_COUNT)
    capture.release()

    if not fps or fps <= 0:
        return 0.0

    return float(total_frames / fps)


def get_frame_at_timestamp(video_path: str, timestamp: float):
    capture = cv2.VideoCapture(video_path)
    fps = capture.get(cv2.CAP_PROP_FPS)

    if not fps or fps <= 0:
        fps = 30

    frame_index = max(0, int(timestamp * fps))
    capture.set(cv2.CAP_PROP_POS_FRAMES, frame_index)
    success, frame = capture.read()
    capture.release()

    if not success:
        return None

    return frame


def detect_text(frame: np.ndarray, min_chars: int = 3, ocr_language: str = "eng") -> str:
    """
    Return OCR text from a frame, or an empty string if text is too short.
    Uses stable OCR behavior and supports selected OCR languages.
    """
    gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)

    gray = cv2.resize(gray, None, fx=1.5, fy=1.5, interpolation=cv2.INTER_CUBIC)
    gray = cv2.GaussianBlur(gray, (3, 3), 0)
    _, thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)

    try:
        raw_text = pytesseract.image_to_string(thresh, lang=ocr_language)
    except pytesseract.TesseractError:
        raw_text = pytesseract.image_to_string(thresh, lang="eng")

    cleaned = clean_text(raw_text)

    if len(cleaned) < min_chars:
        return ""

    return cleaned


def frame_difference(a: np.ndarray, b: np.ndarray) -> float:
    a_small = cv2.resize(a, (160, 90))
    b_small = cv2.resize(b, (160, 90))
    diff = cv2.absdiff(a_small, b_small)
    return float(np.mean(diff))


def text_similarity(a: str, b: str) -> float:
    return SequenceMatcher(None, a.lower().strip(), b.lower().strip()).ratio()


def save_frame(frame: np.ndarray, output_dir: Path, index: int, timestamp: float) -> str:
    frame_filename = output_dir / f"frame_{index:03d}_{timestamp:.2f}s.jpg"
    cv2.imwrite(str(frame_filename), frame)
    return str(frame_filename)


def choose_best_frame_from_sequence(sequence: List[Dict]) -> Dict:
    return max(sequence, key=lambda item: (len(item["text"]), item["timestamp"]))


def extract_text_frames(
    video_path: str,
    sample_every_seconds: float = 0.25,
    min_chars: int = 3,
    duplicate_threshold: float = 8.0,
    sequence_gap_seconds: float = 0.25,
    ocr_language: str = "eng",
) -> List[FrameResult]:
    output_dir = Path(tempfile.mkdtemp())
    capture = cv2.VideoCapture(video_path)

    fps = capture.get(cv2.CAP_PROP_FPS)
    if not fps or fps <= 0:
        fps = 30

    frame_interval = max(1, int(fps * sample_every_seconds))
    max_gap_frames = max(1, int(fps * sequence_gap_seconds))

    text_frames: List[FrameResult] = []
    current_sequence: List[Dict] = []
    last_text_frame_index = None
    last_saved_frame = None
    last_saved_text = ""
    frame_index = 0

    def finish_sequence():
        nonlocal last_saved_frame, last_saved_text

        if not current_sequence:
            return

        best = choose_best_frame_from_sequence(current_sequence)

        is_duplicate = False
        if last_saved_frame is not None:
            diff_score = frame_difference(best["frame"], last_saved_frame)
            text_score = text_similarity(best["text"], last_saved_text)

            image_is_similar = diff_score < duplicate_threshold
            text_is_similar = text_score > 0.82
            is_duplicate = image_is_similar and text_is_similar

        if not is_duplicate:
            frame_path = save_frame(
                best["frame"],
                output_dir,
                len(text_frames) + 1,
                best["timestamp"],
            )
            text_frames.append((frame_path, best["timestamp"], best["text"]))
            last_saved_frame = best["frame"].copy()
            last_saved_text = best["text"]

        current_sequence.clear()

    while True:
        success, frame = capture.read()
        if not success:
            break

        if frame_index % frame_interval == 0:
            timestamp = frame_index / fps
            detected_text = detect_text(frame, min_chars=min_chars, ocr_language=ocr_language)

            if detected_text:
                if (
                    last_text_frame_index is not None
                    and frame_index - last_text_frame_index > max_gap_frames
                ):
                    finish_sequence()

                current_sequence.append(
                    {
                        "frame": frame.copy(),
                        "timestamp": timestamp,
                        "text": detected_text,
                    }
                )
                last_text_frame_index = frame_index

            else:
                if (
                    current_sequence
                    and last_text_frame_index is not None
                    and frame_index - last_text_frame_index > max_gap_frames
                ):
                    finish_sequence()

        frame_index += 1

    finish_sequence()
    capture.release()
    return text_frames


def create_powerpoint(text_frames: List[FrameResult], video_path: str) -> str:
    """
    Create a PowerPoint storyboard deck.

    Layout:
    - One captured frame per slide
    - Captured frame on the left
    - Blank notes area on the right
    - Final slide reminds user to reference original motion ad separately
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for i, (frame_path, timestamp, detected_text) in enumerate(text_frames, start=1):
        slide = prs.slides.add_slide(blank_layout)

        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.45))
        title_box.text_frame.text = f"Motion Storyboard QC — Frame {i}"

        timestamp_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.72), Inches(5.8), Inches(0.3))
        timestamp_box.text_frame.text = f"Timestamp: {timestamp:.2f}s"

        slide.shapes.add_picture(
            frame_path,
            Inches(0.6),
            Inches(1.15),
            width=Inches(5.9),
            height=Inches(4.45),
        )

        notes_label = slide.shapes.add_textbox(Inches(6.85), Inches(0.72), Inches(5.85), Inches(0.3))
        notes_label.text_frame.text = "Notes / Copy / QC"

        blank_notes_box = slide.shapes.add_textbox(Inches(6.85), Inches(1.15), Inches(5.85), Inches(5.9))
        blank_notes_box.text_frame.text = ""
        blank_notes_box.line.width = 1

    final_slide = prs.slides.add_slide(blank_layout)

    title_box = final_slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.5))
    title_box.text_frame.text = "Full Motion Reference"

    note_box = final_slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(10.8), Inches(1.5))
    note_box.text_frame.text = (
        "Please reference the original motion ad/video separately during QC review. "
        "This keeps the generated deck lighter and more reliable in Google Slides."
    )

    output_path = str(Path(tempfile.mkdtemp()) / "motion_storyboard_deck.pptx")
    prs.save(output_path)
    return output_path


def combine_and_sort_frames(auto_frames: List[FrameResult], manual_frames: List[FrameResult]) -> List[FrameResult]:
    combined = auto_frames + manual_frames
    combined.sort(key=lambda item: item[1])
    return combined


st.set_page_config(page_title="Motion Storyboard QC", layout="wide")

st.title("Motion Storyboard QC Generator")
st.write("Upload a motion video. The tool finds frames with visible copy and turns them into a storyboard deck.")

uploaded_video = st.file_uploader("Upload video", type=["mp4", "mov", "m4v", "avi"])

with st.sidebar:
    st.header("OCR Language")
    selected_language_label = st.selectbox(
        "Choose the main text language",
        list(OCR_LANGUAGES.keys()),
        index=0,
    )
    selected_language_code = OCR_LANGUAGES[selected_language_label]

    st.caption(
        "If the ad has English plus another language, choose a combined option like English + Korean."
    )

    st.header("QC Detection Settings")
    sample_rate = st.slider("Scan every X seconds", 0.10, 2.0, 0.25, 0.05)
    min_chars = st.slider("Minimum OCR characters", 1, 20, 3, 1)
    duplicate_threshold = st.slider("Duplicate filter strength", 1.0, 25.0, 8.0, 1.0)
    sequence_gap = st.slider("Moving text grouping window", 0.25, 2.0, 0.25, 0.25)

st.info(
    "Stable default: scan every 0.25s, minimum OCR characters 3, duplicate strength 8, moving text window 0.25. "
    "Use Manual Add Frame if the automatic scan misses a moment."
)

if "auto_frames" not in st.session_state:
    st.session_state.auto_frames = []

if "manual_frames" not in st.session_state:
    st.session_state.manual_frames = []

if uploaded_video:
    with tempfile.NamedTemporaryFile(delete=False, suffix=Path(uploaded_video.name).suffix) as tmp:
        tmp.write(uploaded_video.read())
        video_path = tmp.name

    if st.session_state.get("current_video_name") != uploaded_video.name:
        st.session_state.current_video_name = uploaded_video.name
        st.session_state.auto_frames = []
        st.session_state.manual_frames = []

    duration = get_video_duration(video_path)

    st.video(video_path)

    if st.button("Generate Automatic Storyboard Frames"):
        with st.spinner(f"Scanning video for text frames using {selected_language_label} OCR..."):
            st.session_state.auto_frames = extract_text_frames(
                video_path,
                sample_every_seconds=sample_rate,
                min_chars=min_chars,
                duplicate_threshold=duplicate_threshold,
                sequence_gap_seconds=sequence_gap,
                ocr_language=selected_language_code,
            )

        st.success(f"Found {len(st.session_state.auto_frames)} automatic storyboard frames with text.")

    st.divider()

    st.subheader("Manual Add Frame")
    st.write(
        "If the automatic scan misses a frame, use the timestamp slider below, preview the frame, then click Add This Frame."
    )

    manual_timestamp = st.slider(
        "Choose timestamp to capture manually",
        min_value=0.0,
        max_value=max(duration, 0.1),
        value=0.0,
        step=0.05,
    )

    preview_frame = get_frame_at_timestamp(video_path, manual_timestamp)

    if preview_frame is not None:
        preview_rgb = cv2.cvtColor(preview_frame, cv2.COLOR_BGR2RGB)
        st.image(preview_rgb, caption=f"Manual preview at {manual_timestamp:.2f}s", width=500)

        manual_note = st.text_input(
            "Optional label for this manual frame",
            value="Manually added frame",
        )

        if st.button("Add This Frame"):
            manual_dir = Path(tempfile.mkdtemp())
            frame_path = save_frame(
                preview_frame,
                manual_dir,
                len(st.session_state.manual_frames) + 1,
                manual_timestamp,
            )

            st.session_state.manual_frames.append(
                (frame_path, manual_timestamp, manual_note)
            )
            st.success(f"Added manual frame at {manual_timestamp:.2f}s.")

    st.divider()

    combined_frames = combine_and_sort_frames(
        st.session_state.auto_frames,
        st.session_state.manual_frames,
    )

    st.subheader("Storyboard Frames to Export")
    st.write(
        f"Automatic frames: {len(st.session_state.auto_frames)} | Manual frames: {len(st.session_state.manual_frames)} | Total: {len(combined_frames)}"
    )

    if combined_frames:
        cols = st.columns(3)
        for i, (frame_path, timestamp, detected_text) in enumerate(combined_frames):
            with cols[i % 3]:
                st.image(frame_path, caption=f"{timestamp:.2f}s")
                st.caption(detected_text[:200])

        if st.button("Clear Manual Frames"):
            st.session_state.manual_frames = []
            st.rerun()

        with st.spinner("Building storyboard deck with review slides..."):
            pptx_path = create_powerpoint(combined_frames, video_path)

        with open(pptx_path, "rb") as f:
            st.download_button(
                "Download Storyboard Deck for Google Slides",
                data=f,
                file_name="motion_storyboard_deck.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
    else:
        st.warning("No frames selected yet. Run automatic detection or manually add frames.")
